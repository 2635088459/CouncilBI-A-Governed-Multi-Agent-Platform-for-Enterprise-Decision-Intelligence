"""Unstructured file ingestion: text extraction and semantic-aware chunking.

See FR-FV10-010 in
spec/final-version/en/10-user-file-upload-and-hybrid-analysis.spec.md:
extract plain text from PDF/DOCX/TXT/MD/PPTX, then split it into
300-500 token chunks with a 50 token overlap, each tagged with the
org/user/file scope it must be filtered by at retrieval time.
"""

from __future__ import annotations

import io
import re
from dataclasses import dataclass
from typing import Any, cast

import pypdf
from docx import Document
from pptx import Presentation


_SENTENCE_BOUNDARY_PATTERN = re.compile(r"(?<=[.!?。！？])\s+")
_TOKEN_PATTERN = re.compile(r"\S+")

DEFAULT_MAX_TOKENS = 400
DEFAULT_OVERLAP_TOKENS = 50


class TextExtractor:
    """Extract plain text from every unstructured format in the FR-FV10-002 whitelist."""

    def extract(self, content_bytes: bytes, extension: str) -> str:
        normalized_extension = extension.lower().lstrip(".")
        if normalized_extension == "pdf":
            return self.extract_pdf(content_bytes)
        if normalized_extension == "docx":
            return self.extract_docx(content_bytes)
        if normalized_extension == "pptx":
            return self.extract_pptx(content_bytes)
        if normalized_extension in ("txt", "md"):
            return content_bytes.decode("utf-8")
        raise ValueError(f"unsupported unstructured file extension: {extension}")

    def extract_pdf(self, content_bytes: bytes) -> str:
        reader = pypdf.PdfReader(io.BytesIO(content_bytes))
        page_texts = [page.extract_text() or "" for page in reader.pages]
        return "\n".join(page_texts).strip()

    def extract_docx(self, content_bytes: bytes) -> str:
        document = Document(io.BytesIO(content_bytes))
        paragraph_texts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        return "\n".join(paragraph_texts).strip()

    def extract_pptx(self, content_bytes: bytes) -> str:
        presentation = Presentation(io.BytesIO(content_bytes))
        slide_texts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = cast(Any, shape).text_frame.text
                if text.strip():
                    slide_texts.append(text)
        return "\n".join(slide_texts).strip()


@dataclass(frozen=True, slots=True)
class TextChunk:
    text: str
    chunk_index: int
    org_id: str = ""
    user_id: str = ""
    file_id: str = ""


class SentenceAwareChunker:
    """Split extracted text into overlapping chunks without breaking a sentence mid-word."""

    def chunk(
        self,
        text: str,
        *,
        max_tokens: int = DEFAULT_MAX_TOKENS,
        overlap: int = DEFAULT_OVERLAP_TOKENS,
        org_id: str = "",
        user_id: str = "",
        file_id: str = "",
    ) -> tuple[TextChunk, ...]:
        sentence_token_lists = [
            _TOKEN_PATTERN.findall(sentence)
            for sentence in _split_sentences(text)
            if sentence.strip()
        ]

        finalized_chunks: list[list[str]] = []
        current_tokens: list[str] = []
        has_new_content = False

        for sentence_tokens in sentence_token_lists:
            prospective_length = len(current_tokens) + len(sentence_tokens)
            if has_new_content and prospective_length > max_tokens:
                finalized_chunks.append(current_tokens)
                current_tokens = current_tokens[-overlap:] if overlap > 0 else []
                has_new_content = False
            current_tokens.extend(sentence_tokens)
            has_new_content = True

        if current_tokens:
            finalized_chunks.append(current_tokens)

        return tuple(
            TextChunk(
                text=" ".join(tokens),
                chunk_index=index,
                org_id=org_id,
                user_id=user_id,
                file_id=file_id,
            )
            for index, tokens in enumerate(finalized_chunks)
        )


def _split_sentences(text: str) -> list[str]:
    return _SENTENCE_BOUNDARY_PATTERN.split(text.strip())
